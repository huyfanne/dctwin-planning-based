"""DCTwin presentation deck toolkit.

Two things:
  1. annotate(): overlay clean numbered markers / highlight boxes on a real screenshot
     (legend text lives on the slide, keeping the image uncluttered).
  2. python-pptx helpers with a dark theme matching the DCTwin app, including speaker notes.
"""
from __future__ import annotations
import os
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- theme (matches the app's dark navy / cyan look) ----
BG      = RGBColor(0x0A, 0x0E, 0x1A)   # near-black navy
PANEL   = RGBColor(0x12, 0x18, 0x28)   # panel
ACCENT  = RGBColor(0x36, 0xE0, 0xE0)   # cyan
ACCENT2 = RGBColor(0x66, 0xD9, 0x9A)   # mint/green (safe)
WARN    = RGBColor(0xF0, 0xA8, 0x3A)   # amber
DANGER  = RGBColor(0xFF, 0x6B, 0x6B)   # red
TEXT    = RGBColor(0xE8, 0xEE, 0xF5)   # near-white
DIM     = RGBColor(0x97, 0xA6, 0xBC)   # dim grey-blue
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

_FONTDIR = None
def _font(size, bold=False):
    global _FONTDIR
    if _FONTDIR is None:
        import matplotlib
        _FONTDIR = os.path.join(os.path.dirname(matplotlib.__file__), "mpl-data", "fonts", "ttf")
    name = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    return ImageFont.truetype(os.path.join(_FONTDIR, name), size)

# RGBA accent for PIL
PIL_ACCENT = (0x36, 0xE0, 0xE0)
PIL_AMBER  = (0xF0, 0xA8, 0x3A)
PIL_RED    = (0xFF, 0x6B, 0x6B)
PIL_GREEN  = (0x66, 0xD9, 0x9A)


def annotate(src, dst, markers=None, boxes=None, scale_to=1800, crop=None):
    """markers: list of dict(n, x, y, color?) with x,y in [0,1] relative coords.
       boxes:   list of dict(x0,y0,x1,y1, color?, label?) relative coords.
       crop:    optional (x0,y0,x1,y1) relative box applied FIRST; markers/boxes
                coords are then relative to the cropped image.
       Draws on a copy; downscales width to scale_to for deck size."""
    im = Image.open(src).convert("RGB")
    if crop:
        W0, H0 = im.size
        im = im.crop((int(crop[0]*W0), int(crop[1]*H0), int(crop[2]*W0), int(crop[3]*H0)))
    W, H = im.size
    draw = ImageDraw.Draw(im, "RGBA")
    r = max(16, int(W * 0.016))               # marker radius
    lw = max(3, int(W * 0.0025))              # line width

    for b in (boxes or []):
        col = b.get("color", PIL_ACCENT)
        x0, y0, x1, y1 = int(b["x0"]*W), int(b["y0"]*H), int(b["x1"]*W), int(b["y1"]*H)
        draw.rectangle([x0, y0, x1, y1], outline=col + (255,), width=lw)
        if b.get("label"):
            f = _font(int(r*1.1), bold=True)
            draw.text((x0+6, max(0, y0 - int(r*1.4))), b["label"], fill=col + (255,), font=f)

    f = _font(int(r*1.15), bold=True)
    for m in (markers or []):
        col = m.get("color", PIL_ACCENT)
        cx, cy = int(m["x"]*W), int(m["y"]*H)
        # halo + filled circle + number
        draw.ellipse([cx-r-lw, cy-r-lw, cx+r+lw, cy+r+lw], fill=(0, 0, 0, 170))
        draw.ellipse([cx-r, cy-r, cx+r, cy+r], fill=col + (255,), outline=(255, 255, 255, 255), width=lw)
        t = str(m["n"])
        bb = draw.textbbox((0, 0), t, font=f)
        draw.text((cx-(bb[2]-bb[0])/2, cy-(bb[3]-bb[1])/2 - bb[1]), t, fill=(8, 12, 22, 255), font=f)

    if scale_to and W > scale_to:
        h = int(H * scale_to / W)
        im = im.resize((scale_to, h), Image.LANCZOS)
    im.save(dst)
    return dst


# ---------------- pptx helpers ----------------
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

def new_deck():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs

def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)  # send to back
    return s

def _tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tb, tf

def _set(p, text, size, color, bold=False, italic=False):
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size); run.font.color.rgb = color
        run.font.bold = bold; run.font.italic = italic
        run.font.name = "Calibri"

def _bar(slide, prs, color=ACCENT, x=Inches(0.0), y=Inches(0.0), w=None, h=Inches(0.09)):
    w = w or prs.slide_width
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background(); r.shadow.inherit = False
    return r

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def title_slide(prs, title, subtitle, thesis, footer=""):
    s = _blank(prs)
    _bar(s, prs, ACCENT, y=Inches(2.55), x=Inches(1.0), w=Inches(2.2), h=Inches(0.06))
    _, tf = _tb(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(2.0))
    _set(tf.paragraphs[0], title, 40, TEXT, bold=True)
    p = tf.add_paragraph(); _set(p, subtitle, 21, ACCENT); p.space_before = Pt(10)
    _, tf2 = _tb(s, Inches(1.0), Inches(4.9), Inches(11.0), Inches(1.6))
    _set(tf2.paragraphs[0], thesis, 16, DIM, italic=True)
    if footer:
        _, tff = _tb(s, Inches(1.0), Inches(6.7), Inches(11.0), Inches(0.5))
        _set(tff.paragraphs[0], footer, 12, DIM)
    return s

def section_slide(prs, act, title, subtitle):
    s = _blank(prs)
    _, tfk = _tb(s, Inches(1.0), Inches(2.7), Inches(11), Inches(0.6))
    _set(tfk.paragraphs[0], act, 15, ACCENT, bold=True)
    _, tf = _tb(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(1.4))
    _set(tf.paragraphs[0], title, 33, TEXT, bold=True)
    _, tf2 = _tb(s, Inches(1.0), Inches(4.7), Inches(11.0), Inches(1.4))
    _set(tf2.paragraphs[0], subtitle, 17, DIM)
    _bar(s, prs, ACCENT, x=Inches(1.0), y=Inches(2.55), w=Inches(1.6), h=Inches(0.05))
    return s

def _header(slide, prs, kicker, title):
    _bar(slide, prs, ACCENT, h=Inches(0.08))
    _, tfk = _tb(slide, Inches(0.55), Inches(0.22), Inches(12), Inches(0.4))
    _set(tfk.paragraphs[0], kicker, 13, ACCENT, bold=True)
    _, tf = _tb(slide, Inches(0.55), Inches(0.56), Inches(12.2), Inches(0.7))
    _set(tf.paragraphs[0], title, 25, TEXT, bold=True)

def image_slide(prs, kicker, title, image_path, legend, note_text, img_frac=0.62):
    """image on the left, numbered legend bullets on the right. legend: list of (n_or_None, text)."""
    s = _blank(prs); _header(s, prs, kicker, title)
    # image area
    top = Inches(1.45); left = Inches(0.55)
    maxw = prs.slide_width * img_frac - left
    maxh = prs.slide_height - top - Inches(0.45)
    im = Image.open(image_path); iw, ih = im.size
    scale = min(maxw / iw, maxh / ih)
    w = int(iw * scale); h = int(ih * scale)
    # frame
    fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - Emu(40000), top - Emu(40000), Emu(w) + Emu(80000), Emu(h) + Emu(80000))
    fr.fill.solid(); fr.fill.fore_color.rgb = PANEL; fr.line.color.rgb = ACCENT; fr.line.width = Pt(1); fr.shadow.inherit = False
    s.shapes.add_picture(image_path, left, top, width=Emu(w), height=Emu(h))
    # legend
    lx = left + Emu(w) + Inches(0.35)
    lw = prs.slide_width - lx - Inches(0.4)
    _, tf = _tb(s, lx, top, lw, maxh)
    first = True
    for item in legend:
        n, text = item
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        if n is None:
            _set(p, text, 14, ACCENT2, bold=True); p.space_before = Pt(8)
        else:
            _set(p, f"{n}", 14, BG, bold=True)
            p.runs[0].font.color.rgb = ACCENT
            r2 = p.add_run(); r2.text = "  " + text
            r2.font.size = Pt(13.5); r2.font.color.rgb = TEXT; r2.font.name = "Calibri"
            p.space_before = Pt(7)
    notes(s, note_text)
    return s

def bullets_slide(prs, kicker, title, bullets, note_text):
    """bullets: list of (level, text) ; level 0 = headline accent, 1 = normal, 2 = sub-dim."""
    s = _blank(prs); _header(s, prs, kicker, title)
    _, tf = _tb(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(5.5))
    first = True
    for lvl, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        if lvl == 0:
            _set(p, text, 18, ACCENT2, bold=True); p.space_before = Pt(12)
        elif lvl == 1:
            _set(p, "•  " + text, 15.5, TEXT); p.space_before = Pt(6)
        else:
            _set(p, "    – " + text, 13.5, DIM); p.space_before = Pt(3)
    notes(s, note_text)
    return s

def qa_slide(prs, kicker, title, qas, note_text):
    s = _blank(prs); _header(s, prs, kicker, title)
    _, tf = _tb(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
    first = True
    for q, a in qas:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        _set(p, "Q  " + q, 14.5, ACCENT, bold=True); p.space_before = Pt(11)
        pa = tf.add_paragraph(); _set(pa, "A  " + a, 13, TEXT); pa.space_before = Pt(2)
    notes(s, note_text)
    return s

def save(prs, path):
    prs.save(path)
    return path
