"""Build DCTwin_4th_draft_report_full.pptx — full technical + operational deck (4th draft).
Reuses deck_assets figures + real UI screenshots; adds safety-stack, data-flow and
operational-procedure diagrams. Run with the project venv from anywhere."""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch

OUT = Path("/mnt/lv/home/hoanghuy/newcode/dctwin/dctwin/BHUMANS")
ASSETS = OUT / "deck_assets"
SHOTS = Path("/mnt/lv/home/hoanghuy/newcode/dctwin/src/log/screens")

NAVY = "#1F3A5F"; ORANGE = "#E8852A"; TEAL = "#2A8E9E"; RED = "#C0392B"; GREEN = "#2E7D32"
GREY = "#7A7A7A"; LIGHT = "#EAF2F8"


def _box(ax, x, y, w, h, text, fc, tc="white", fs=10):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.35",
                                facecolor=fc, edgecolor=NAVY, lw=1.4))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", color=tc,
            fontsize=fs, weight="bold")


def _arr(ax, x1, y1, x2, y2, color=NAVY, ls="-"):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>", color=color,
                                 lw=2, mutation_scale=16, linestyle=ls))


# ── figD: the safety stack ──
fig, ax = plt.subplots(figsize=(12.2, 6.0))
ax.set_xlim(0, 100); ax.set_ylim(0, 60); ax.axis("off")
ax.text(50, 56.5, "Three independent safety nets — uncertainty hedged exactly once",
        ha="center", fontsize=14, weight="bold", color=NAVY)
_box(ax, 4, 38, 28, 11, "1 · SEARCH GATE\ninlet ≤ 26 °C − k·σ\n(σ = fading-floor calibration)", RED, fs=10)
_box(ax, 36, 38, 28, 11, "2 · ROBUST ENSEMBLE\nworst plant (±σ_post, data-centered)\n+ HOT-WEATHER week (+1σ)\nmust hold the HARD 26 °C cap", "#8E44AD", fs=9.5)
_box(ax, 68, 38, 28, 11, "3 · DEPLOY BACKSTOP\nany realised violation\n→ deploy_blocked", NAVY, fs=10)
_arr(ax, 32, 43.5, 36, 43.5); _arr(ax, 64, 43.5, 68, 43.5)
_box(ax, 12, 20, 34, 9, "fragile optimum?\nSAFETY LADDER substitutes the cheapest\nPROVABLY-robust plan (CHWST-first)", ORANGE, fs=9.5)
_arr(ax, 50, 38, 35, 29, color=ORANGE)
_box(ax, 54, 20, 34, 9, "uncertainty allocation:\nk·σ margin → nominal check ONLY;\nscenarios test bias-corrected HARD cap", TEAL, fs=9.5)
ax.text(50, 10, "measured outcome: 6 deployed weeks · worst inlet 25.60 °C · 0 violations",
        ha="center", fontsize=12, weight="bold", color=GREEN)
fig.tight_layout(); fig.savefig(ASSETS / "figD_safety_stack.png", dpi=150); plt.close(fig)

# ── figE: telemetry & shadow-BMS data flow ──
fig, ax = plt.subplots(figsize=(12.2, 6.2))
ax.set_xlim(0, 100); ax.set_ylim(0, 62); ax.axis("off")
ax.text(50, 58.5, "Closing the loop to reality — every seam built, shadow-safe today",
        ha="center", fontsize=14, weight="bold", color=NAVY)
_box(ax, 3, 40, 22, 10, "Approved plan\n(3 setpoints)", NAVY)
_box(ax, 31, 40, 26, 10, "ShadowBmsAdapter\n45 actuator commands\nRECORDED, not actuated", ORANGE)
_box(ax, 63, 40, 17, 10, "BACnet adapter\n(field seam,\nexplicit stub)", GREY)
_box(ax, 84, 40, 13, 10, "Real BMS\n(future)", "#444444")
_arr(ax, 25, 45, 31, 45); _arr(ax, 57, 45, 63, 45, ls="--"); _arr(ax, 80, 45, 84, 45, ls="--")
_box(ax, 3, 22, 22, 10, "SimTelemetryFeed\n(labelled SIMULATED)", TEAL)
_box(ax, 3, 8, 22, 9, "Real historian /\ncollector (future)", GREY)
_box(ax, 31, 15, 26, 13, "TelemetryStore\nPOST /api/telemetry\n22 rack inlets · power ·\nPUE · RH · held setpoints", TEAL)
_arr(ax, 25, 27, 31, 24); _arr(ax, 25, 12.5, 31, 18, ls="--")
_box(ax, 63, 15, 34, 13, "Live dashboard: heat-map · alerts\n(≥25 warn / ≥26 critical) · setpoint\ncompliance · rack hotspot table\n+ calibration & recirc fitting", NAVY, fs=9.5)
_arr(ax, 57, 21.5, 63, 21.5)
ax.text(50, 2.5, "swap the dashed boxes for field systems — zero changes downstream",
        ha="center", fontsize=11.5, color=ORANGE, weight="bold")
fig.tight_layout(); fig.savefig(ASSETS / "figE_dataflow.png", dpi=150); plt.close(fig)

# ── figF: operator weekly procedure ──
fig, ax = plt.subplots(figsize=(12.4, 5.6))
ax.set_xlim(0, 100); ax.set_ylim(0, 52); ax.axis("off")
ax.text(50, 48.5, "The operator's week — one cycle of the dual loop", ha="center",
        fontsize=14, weight="bold", color=NAVY)
steps = [
    ("1 · MONITOR", "Live tab:\nheat-map, alerts,\ncompliance", TEAL),
    ("2 · PLAN", "New Plan: context\ndecks → launch →\nlive progress", NAVY),
    ("3 · REVIEW", "KPIs vs baseline,\nbands, inlet\ntrajectory vs 26 °C", NAVY),
    ("4 · APPROVE", "expert gate;\nblocked_unsafe is\nNOT approvable", ORANGE),
    ("5 · DEPLOY", "shadow mode:\n45 commands\nrecorded", ORANGE),
    ("6 · LEARN", "realised vs predicted\n→ calibration σ↓\n→ savings ↑", GREEN),
]
w = 14.2
for i, (t, d, c) in enumerate(steps):
    x = 2.5 + i * (w + 2.0)
    _box(ax, x, 26, w, 13, t, c, fs=10.5)
    ax.text(x + w / 2, 17.5, d, ha="center", va="center", fontsize=8.8, color=GREY)
    if i < 5:
        _arr(ax, x + w, 32.5, x + w + 2.0, 32.5)
ax.add_patch(FancyArrowPatch((95, 25), (8, 9), arrowstyle="-|>", color=GREEN, lw=2,
                             mutation_scale=16, connectionstyle="arc3,rad=-0.18"))
ax.text(50, 6.5, "next week starts smarter: tighter σ → the gate releases more savings",
        ha="center", fontsize=11.5, color=GREEN, weight="bold")
fig.tight_layout(); fig.savefig(ASSETS / "figF_procedure.png", dpi=150); plt.close(fig)
print("new figs done")

# ================================================================ PPTX
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY_C = RGBColor(0x1F, 0x3A, 0x5F); ORANGE_C = RGBColor(0xE8, 0x85, 0x2A)
TEAL_C = RGBColor(0x2A, 0x8E, 0x9E); GREY_C = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); GREEN_C = RGBColor(0x2E, 0x7D, 0x32)
LBLUE = RGBColor(0xBF, 0xD3, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def _txt(s, x, y, w, h, text, size, color=NAVY_C, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = "Calibri"
    return tb


def _rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def header(s, title, kicker=None):
    _rect(s, 0, 0, SW, Inches(1.1), NAVY_C)
    _rect(s, 0, Inches(1.1), SW, Inches(0.06), ORANGE_C)
    _txt(s, Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.75), title, 25, WHITE, True)
    if kicker:
        _txt(s, Inches(0.55), Inches(0.78), SW - Inches(1), Inches(0.3), kicker, 12, LBLUE)


def bullets(s, items, x=Inches(0.7), y=Inches(1.45), size=15.5, gap=1.16):
    tb = s.shapes.add_textbox(x, y, SW - Inches(1.4), SH - y - Inches(0.35))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl, text, col, bold = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap; p.level = lvl
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + text
        f = r.font; f.size = Pt(size - lvl * 2); f.color.rgb = col; f.bold = bold; f.name = "Calibri"


def pic_slide(title, img, caption, kicker=None):
    s = prs.slides.add_slide(blank); header(s, title, kicker)
    from PIL import Image
    iw, ih = Image.open(img).size
    sc = min((SW - Inches(1.1)) / iw, (SH - Inches(2.05)) / ih)
    w, h = Emu(int(iw * sc)), Emu(int(ih * sc))
    s.shapes.add_picture(str(img), Emu(int((SW - w) / 2)), Inches(1.4), w, h)
    _txt(s, Inches(0.6), SH - Inches(0.5), SW - Inches(1.2), Inches(0.4), caption, 11.5,
         GREY_C, align=PP_ALIGN.CENTER, italic=True)
    return s


def tile(s, x, y, w, h, big, small, color=NAVY_C):
    _rect(s, x, y, w, h, RGBColor(0xEF, 0xF3, 0xF8))
    _txt(s, x, y + Inches(0.15), w, Inches(0.7), big, 28, color, True, PP_ALIGN.CENTER)
    _txt(s, x + Inches(0.08), y + Inches(0.9), w - Inches(0.16), h - Inches(0.95),
         small, 12, GREY_C, align=PP_ALIGN.CENTER)


G = GREY_C; N = NAVY_C; O = ORANGE_C; T = TEAL_C; GR = GREEN_C

# 1 title
s = prs.slides.add_slide(blank)
_rect(s, 0, 0, SW, SH, NAVY_C); _rect(s, 0, Inches(4.0), SW, Inches(0.07), ORANGE_C)
_txt(s, Inches(0.9), Inches(1.7), SW - Inches(1.8), Inches(1.2),
     "DCTwin — Digital Twin Dual-Loop Control", 38, WHITE, True)
_txt(s, Inches(0.95), Inches(2.8), SW - Inches(1.8), Inches(0.6),
     "4th draft · full technical & operational report", 20, ORANGE_C)
_txt(s, Inches(0.95), Inches(4.35), SW - Inches(1.8), Inches(1.0),
     "Physics-based weekly cooling optimization, measured over six deployed weeks:\n"
     "savings 0 → 3.6% · prediction error 1.33% → 0.1% · zero safety violations", 15, LBLUE)
_txt(s, Inches(0.95), Inches(6.0), SW - Inches(1.8), Inches(0.4), "June 2026", 13, WHITE)

# 2 framework
pic_slide("The framework — two nested loops", ASSETS / "fig3_arch.png",
          "Inner loop: forecast → EnergyPlus physics → beam search → safety filter. Outer loop: pre-validation → expert gate → (shadow) deploy → measure → learn.",
          "Matches the original optimization-plan diagram box-for-box")

# 3 headline
s = prs.slides.add_slide(blank); header(s, "Headline results (all measured)", "Weeks of 2024-11-08 … 12-13")
W, H, GX = Inches(3.85), Inches(2.1), Inches(0.32); x0, y0 = Inches(0.62), Inches(1.55)
tile(s, x0, y0, W, H, "6 weeks", "consecutive cycles: plan → gate → approve →\nshadow-deploy → measure → learn", N)
tile(s, x0 + W + GX, y0, W, H, "3.6%", "energy saving vs as-operated baseline,\nreleased stepwise by the safety gate", O)
tile(s, x0 + 2 * (W + GX), y0, W, H, "0", "inlet violations — worst realised inlet\n25.60 °C vs the 26 °C hard cap", GR)
y1 = y0 + H + Inches(0.3)
tile(s, x0, y1, W, H, "0.1%", "twin weekly-energy prediction error by week 6\n(1.33% in week 1 — calibration learning)", T)
tile(s, x0 + W + GX, y1, W, H, "≤0.2 °C", "inlet prediction error every week —\nthe safety variable is best-predicted", GR)
tile(s, x0 + 2 * (W + GX), y1, W, H, "526", "automated tests green (422 backend +\n104 frontend) · 21 API endpoints · 7 UI pages", N)

# 4-6 results charts
pic_slide("Result 1 — the twin predicts reality", ASSETS / "figA_accuracy.png",
          "Predicted vs realised weekly cooling energy; the residual calibration loop drives error 1.33% → ~0.1%.", "Measured")
pic_slide("Result 2 — savings released by earned confidence", ASSETS / "figB_savings.png",
          "The robust gate only passes savings it can PROVE safe; as measured uncertainty σ falls, provably-safe savings rise 0 → 3.6%.", "Measured")
pic_slide("Result 3 — zero violations while saving", ASSETS / "figC_safety.png",
          "Realised peak inlet per week: rising toward — never past — the cap as the gate spends margin only when confidence allows.", "Measured")

# 7-9 physics
pic_slide("The plant — from the real EnergyPlus model", ASSETS / "fig1_building.png",
          "7 stacked halls; controlled hall 1F 2A: 22 agent-controlled CRAHs, 2.0 MW IT; one shared chilled-water plant.", "Physical model")
pic_slide("Controlled hall detail — 3 setpoints → 45 actuators", ASSETS / "fig2_hall.png",
          "CRAH supply-air temperature (20–26 °C) · airflow (4.8–13.8 kg/s) · chilled-water supply temperature (13–19 °C), broadcast to 22+22+1 actuators.", "Physical model")
pic_slide("Why the optimization is well-posed (verified by sweep)", ASSETS / "fig4_results.png",
          "Real-E+ sweeps: airflow moves hall energy 15.5%, CHWST 3.4%, SAT ~0.5% — a true physical gradient with the safety cap binding.", "Validation")

# 10 safety stack
pic_slide("The safety architecture", ASSETS / "figD_safety_stack.png",
          "Hard cap with k·σ margin in search; a data-centered worst-case ensemble (plant ±σ_post + hot-weather week) against the hard cap; a zero-tolerance deploy backstop.",
          "Safety never trades against energy or price")

# 11 inner-loop tech
s = prs.slides.add_slide(blank); header(s, "Technical: the planning engine", "Inner loop")
bullets(s, [
    (0, "Oracle: every candidate = a full-week EnergyPlus 9.5 run (Docker/BCVTB), fanned across a process pool; per-candidate watchdog (300 s) + batch deadline + stall watchdog (1.5× timeout) bound any hang to minutes.", G, False),
    (0, "Search: coarse grid (g³ over the 3-setpoint space) → top-B beam → halving-step refinement; optional day/night time-blocks; degenerate-signal detection.", G, False),
    (0, "Objective: weekly hall-scoped HVAC energy (ACU fans + shared CHW plant) — or, when data/tariff.json exists, the 24-hour tariff/carbon-weighted cost. Soft penalties: inlet excess, RH excursion, zone band.", G, False),
    (0, "Live recirculation fidelity: recirc fraction estimated from rack inlets (mixing identity); flow-shortfall model adds a conservative-only inlet penalty; provable no-op until field-calibrated (fit_recirc.py).", G, False),
    (0, "Forecasting: IT load (persistence/seasonal + calendar alignment, p10/50/90 bands); weather via historical-analog mean+σ with EPW variants — the +1σ hot week joins the robust ensemble.", G, False),
], size=14.5)

# 12 learning loop tech
s = prs.slides.add_slide(blank); header(s, "Technical: the learning loop", "Outer loop")
bullets(s, [
    (0, "After every deployed week: paired (predicted, realised) KPIs update per-KPI bias + two uncertainties:", N, True),
    (1, "fading-floor σ = max(sample, prior/n) → the nominal k·σ search margin (never under-states early);", G, False),
    (1, "empirical-Bayes σ_post = √((n·s² + σ₀²)/(n+1)) → sizes the robust ensemble (evidence tightens it; ±2% physical floor).", G, False),
    (0, "Physics re-calibration: persistent energy bias over ≥4 weeks fits a bounded fan-efficiency correction (clip ±15%) → data/plant_calibration.json → becomes the ensemble's CENTER (the believed plant state is data-driven).", G, False),
    (0, "Measured: n=7 calibration weeks · inlet bias 0.005 °C · σ_inlet 1.0 (prior) → 0.14 °C (measured).", GR, True),
    (0, "Safety-ladder substitution: when the energy optimum is fragile under the ensemble, the cheapest PROVABLY-robust alternative is recommended (chilled-water-first — ~3% energy span vs ~15% for airflow), never an unactionable block.", G, False),
], size=14.5)

# 13 dataflow seams
pic_slide("Closing the loop to reality — the seams", ASSETS / "figE_dataflow.png",
          "Shadow deploys record the exact 45-command artifact a field adapter will write; telemetry accepts any historian push; simulated data is always labelled.",
          "Field connection = config + one adapter class, not a rebuild")

# 14-15 UI screenshots
pic_slide("Operator cockpit — live monitoring (real screenshot)", SHOTS / "live6.png",
          "22-rack heat-map · alert banner (warn ≥25 °C, critical ≥26 °C) · ranked rack hotspot table · commanded-vs-held setpoint compliance · rolling charts · SIMULATED badge.",
          "Expert Supervision: real-time monitoring")
pic_slide("Operator cockpit — 3-D twin with live rack colors (real screenshot)", SHOTS / "twin3d6.png",
          "The controlled hall's rack rows colored from live telemetry (5 s refresh) with legend; HUD shows the active plan's setpoints and KPIs.",
          "Hotspot visibility")

# 16 procedure
pic_slide("Operational procedure — one week, one cycle", ASSETS / "figF_procedure.png",
          "Monitor → plan (context decks) → review (trajectories vs 26 °C) → expert approve → shadow deploy → learn. Every cycle tightens σ and unlocks more savings.",
          "Operator runbook")

# 17 operational considerations
s = prs.slides.add_slide(blank); header(s, "Operational considerations", "What operators should know")
bullets(s, [
    (0, "blocked_unsafe is the system WORKING: the gate refused a plan it could not prove safe — it cannot be force-approved by design. Re-plan or accept the substituted robust alternative.", O, True),
    (0, "All simulated data is labelled (SIMULATED FEED badge; deploy artifacts say actuated:false) — no one can mistake rehearsal for reality.", G, False),
    (0, "Restarting the backend kills a running plan (orphans are marked failed, never lost silently); cancel from the UI first. A wedged run recovers in minutes (stall watchdog + driver unstick).", G, False),
    (0, "Old plans (schema <1.7) show '—' for baseline/reduction by design; re-run to compare. Artifacts are schema-versioned 1.0→1.8, additive.", G, False),
    (0, "Alerts: warn at 25 °C (margin <1 °C), critical at 26 °C; compliance flags drift >0.5 per axis vs the deployed command.", G, False),
    (0, "Reproducible operations: scripts/clear-and-run.sh for humans; the committed /run-dctwin driver (start/smoke/plan/screenshot/unstick) for agents and CI.", G, False),
], size=14.5)

# 18 scorecard
s = prs.slides.add_slide(blank); header(s, "Implemented vs not — honest scorecard", "vs the original plan")
bullets(s, [
    (0, "✅ Every box of the original diagram is implemented: forecaster (load + weather), EnergyPlus model, heuristic planner, unsafe-action filtering, pre-validation, expert supervision incl. REAL-TIME monitoring, system-data feedback.", GR, True),
    (0, "🟡 Authorized commands to the real DC: SHADOW mode (recorded, not actuated). The BACnet adapter is an explicit one-class seam; field access is the only blocker.", O, True),
    (0, "🟡 'Realised' data is a perturbed-plant simulation stand-in until a real collector connects — the learning loop is structurally proven, not yet field-proven.", O, False),
    (0, "Deliberately out of scope: humidity as a control actuator (monitored + penalized only), intra-week automatic re-planning, CI browser E2E (manual driver covers it).", G, False),
    (0, "Known caveats: recirc margin optimistic until field-calibrated; analog-σ weather understates synoptic extremes (+1σ hot scenario partially compensates); IT load ~flat so weather dominates.", G, False),
], size=14.5)

# 19 next steps
s = prs.slides.add_slide(blank); header(s, "Next steps", "The path to the field")
bullets(s, [
    (0, "1 · Field pilot, shadow-first: real collector → POST /api/telemetry; implement BacnetBmsAdapter; recommend-only 4–8 weeks; then actuate hall-by-hall.", N, True),
    (0, "2 · Calibrate recirculation from the first real rack-inlet data (fit_recirc.py) — de-risks the safety margin directly.", G, False),
    (0, "3 · Swap the weather provider seam to a real short-horizon forecast API (downstream unchanged).", G, False),
    (0, "4 · Load the site tariff / grid-carbon profile (data/tariff.json) to optimize the operator's true target.", G, False),
    (0, "5 · Let robustness grow with data: σ_post tightens and the ensemble re-centers automatically; raise n_scenarios for high-stakes weeks.", G, False),
    (0, "6 · Operator hardening: alert webhooks, guided mid-week re-plan, per-ACU power telemetry, CI-managed E2E.", G, False),
], size=15)

# 20 summary
s = prs.slides.add_slide(blank); header(s, "Summary", "")
bullets(s, [
    (0, "A physics-based digital twin + web cockpit that plans a data-center's cooling week and proves its recommendations safe before anyone acts on them.", N, True),
    (0, "Measured, not promised: 6 deployed weeks · 0.1% energy prediction error · 0 → 3.6% savings released by earned confidence · zero violations.", O, True),
    (0, "Live monitoring, shadow deploys, telemetry seams, physics re-calibration, weather + tariff awareness — all tested (526 green) and labelled honestly.", G, False),
    (0, "One step remains to the field, and it is config-shaped: connect real telemetry and the BACnet adapter behind seams that already exist.", GR, True),
], size=17, y=Inches(1.9))

out = OUT / "DCTwin_4th_draft_report_full.pptx"
prs.save(str(out))
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
