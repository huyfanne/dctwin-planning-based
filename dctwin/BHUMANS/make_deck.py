"""Build the professor-facing PPTX report from the 3rd-draft review.
Renders real GDS topology + architecture + validation figures (matplotlib) and
assembles a 16:9 deck (python-pptx). Run from src/ with PYTHONPATH=src. Temp — delete after."""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyBboxPatch, FancyArrowPatch
from pathlib import Path

from webapp.topology import build_hall_topology

OUT = Path("/mnt/lv/home/hoanghuy/newcode/dctwin/dctwin/BHUMANS")
ASSETS = OUT / "deck_assets"
ASSETS.mkdir(parents=True, exist_ok=True)

NAVY = "#1F3A5F"; ORANGE = "#E8852A"; TEAL = "#2A8E9E"; LIGHT = "#D9E4EE"
RED = "#C0392B"; BLUE = "#2E75B6"; GREY = "#7A7A7A"; GREEN = "#2E7D32"

t = build_hall_topology("models/building.json", "configs/dt/dt.prototxt", "1f 2a")
halls = t["building"]["halls"]
plant = t["building"]["plant"]
crahs = t["crahs"]
rack_rows = t["rack_rows"]
HALL_W, HALL_D, _ = t["hall"]["size"]

# ---------------------------------------------------------------- Fig 1: building stack
fig, ax = plt.subplots(figsize=(11.5, 6.2))
for h in halls:
    z0 = h["z0"]; w = h["size"][0]; ht = h["size"][2]; inf = h["infra"]
    ctrl = h["controlled"]
    ax.add_patch(Rectangle((0, z0), w, ht, facecolor=ORANGE if ctrl else LIGHT,
                           edgecolor=NAVY, lw=1.8, alpha=0.95 if ctrl else 0.9))
    label = (f"{h['code']}   ·   {inf['acuTotal']} ACU"
             f"{' (agent-controlled)' if ctrl else ''}   ·   {inf['itPowerKw']/1000:.1f} MW IT")
    ax.text(w / 2, z0 + ht / 2, label, ha="center", va="center",
            fontsize=10.5 if ctrl else 9.5, color="white" if ctrl else NAVY,
            weight="bold" if ctrl else "normal")
ax.annotate("CONTROLLED HALL\n(the planner sets SAT · airflow · CHWST)",
            xy=(HALL_W, 7 + 1.75), xytext=(HALL_W + 6.5, 10.5),
            fontsize=10, color=ORANGE, weight="bold", ha="center",
            arrowprops=dict(arrowstyle="-|>", color=ORANGE, lw=2))
ax.text(HALL_W / 2, -2.4, f"Shared chilled-water plant:  {plant['chiller']} chiller · "
        f"{plant['coolingTower']} cooling tower · {plant['pumps']} pumps",
        ha="center", fontsize=10, color=TEAL, weight="bold")
ax.set_xlim(-3, 60); ax.set_ylim(-3.5, 26)
ax.set_xlabel("hall width (m)", fontsize=9); ax.set_ylabel("building height (m)", fontsize=9)
ax.set_title("GDS Data Center — 7 stacked halls (from real EnergyPlus zone geometry)",
             fontsize=13, weight="bold", color=NAVY)
ax.set_aspect("equal"); ax.spines[["top", "right"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSETS / "fig1_building.png", dpi=150); plt.close(fig)

# ---------------------------------------------------------------- Fig 2: controlled hall plan
fig, ax = plt.subplots(figsize=(11.5, 6.0))
ax.add_patch(Rectangle((0, 0), HALL_W, HALL_D, fill=False, edgecolor=NAVY, lw=2.2))
ax.text(HALL_W / 2, HALL_D + 1.1, f"Controlled Hall {t['hall']['name'].split('Hall ')[-1]}"
        f"  ({HALL_W:.1f} m × {HALL_D:.1f} m)", ha="center", fontsize=12,
        weight="bold", color=NAVY)
# rack rows (cold/hot aisles) down the middle
for r in rack_rows:
    y = r["pos"][1]; col = BLUE if r["aisle"] == "cold" else RED
    ax.add_patch(Rectangle((HALL_W / 2 - 8, y - 0.5), 16, 1.0, facecolor=col,
                           edgecolor="white", lw=0.6, alpha=0.85))
ax.text(HALL_W / 2, rack_rows[0]["pos"][1] - 1.6, "server racks (cold/hot aisles)",
        ha="center", fontsize=8.5, color=GREY)
# CRAHs along south/north walls
for c in crahs:
    x = c["pos"][0]; y = 0.0 if c["wall"] == "south" else HALL_D
    ax.add_patch(Rectangle((x - 0.9, y - 0.9 if c["wall"] == "south" else y - 0.9),
                           1.8, 1.8, facecolor=TEAL, edgecolor=NAVY, lw=0.8))
ax.text(HALL_W / 2, -2.6, f"{len(crahs)} agent-controlled CRAHs (air handlers) along both walls",
        ha="center", fontsize=9.5, color=TEAL, weight="bold")
# chilled-water plant block + CHW links
px, py, pw, ph = -13.5, HALL_D / 2 - 3.2, 9.5, 6.4
ax.add_patch(FancyBboxPatch((px, py), pw, ph, boxstyle="round,pad=0.15",
                            facecolor="#EAF2F8", edgecolor=TEAL, lw=2))
ax.text(px + pw / 2, py + ph / 2, f"Chilled-water\nplant\n{plant['chiller']} chiller · "
        f"{plant['coolingTower']} tower\n{plant['pumps']} pumps", ha="center", va="center",
        fontsize=9, color=NAVY, weight="bold")
for c in crahs[::3]:
    ax.add_patch(FancyArrowPatch((px + pw, py + ph / 2), (c["pos"][0], c["pos"][1] if c["wall"] == "south" else HALL_D),
                                 arrowstyle="-", color=TEAL, lw=0.6, alpha=0.5))
# airflow arrow
ax.add_patch(FancyArrowPatch((6, 1.5), (6, rack_rows[0]["pos"][1] - 1), arrowstyle="-|>",
                             color=BLUE, lw=2))
ax.text(7.2, 3.5, "cold supply air →\nserver inlets (≤ 26 °C)", fontsize=8.5, color=BLUE)
ax.set_xlim(-15, HALL_W + 2); ax.set_ylim(-4, HALL_D + 2.5)
ax.set_aspect("equal"); ax.axis("off")
ax.set_title("Controlled Hall 1F 2A — 22 CRAHs + chilled-water plant (3 setpoints → 45 actuators)",
             fontsize=12.5, weight="bold", color=NAVY, pad=14)
fig.tight_layout(); fig.savefig(ASSETS / "fig2_hall.png", dpi=150); plt.close(fig)

# ---------------------------------------------------------------- Fig 3: dual-loop architecture
fig, ax = plt.subplots(figsize=(12.5, 6.2))
ax.set_xlim(0, 100); ax.set_ylim(0, 60); ax.axis("off")

def box(x, y, w, h, text, fc, tc="white", fs=10):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.4",
                                facecolor=fc, edgecolor=NAVY, lw=1.5))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", color=tc,
            fontsize=fs, weight="bold")

def arrow(x1, y1, x2, y2, color=NAVY):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle="-|>",
                                 color=color, lw=2, mutation_scale=18))

ax.text(50, 57, "Dual-Loop Control Framework", ha="center", fontsize=15, weight="bold", color=NAVY)
# inner loop (top)
ax.text(50, 50.5, "INNER LOOP  —  planning (all in simulation)", ha="center", fontsize=10,
        color=ORANGE, weight="bold")
box(2, 38, 19, 9, "Forecaster\nIT load + weather", TEAL)
box(27, 38, 21, 9, "EnergyPlus 9.5\nPhysics Digital Twin", NAVY)
box(54, 38, 19, 9, "Planner\nbeam search", ORANGE)
box(79, 38, 19, 9, "Safety filter\ninlet ≤ 26 °C", RED)
arrow(21, 42.5, 27, 42.5); arrow(48, 42.5, 54, 42.5); arrow(73, 42.5, 79, 42.5)
arrow(79, 40, 48, 40, color=GREY)  # filtered candidates back to twin
ax.text(63, 36.5, "unsafe options rejected", ha="center", fontsize=8, color=GREY, style="italic")
# outer loop (bottom)
ax.text(50, 27, "OUTER LOOP  —  deployment & learning", ha="center", fontsize=10,
        color=TEAL, weight="bold")
box(2, 14, 19, 9, "Pre-validation\nnominal + worst-case", "#5B7DB1")
box(27, 14, 21, 9, "Expert Supervision\napprove · edit · deploy", "#5B7DB1")
box(54, 14, 19, 9, "Deploy to DC\n(BMS adapter seam)", GREEN)
box(79, 14, 19, 9, "Calibrate\nlearn from realized", "#5B7DB1")
arrow(21, 18.5, 27, 18.5); arrow(48, 18.5, 54, 18.5); arrow(73, 18.5, 79, 18.5)
# connect loops
arrow(89, 38, 89, 23, color=ORANGE)  # plan -> prevalidation side
ax.text(91, 30.5, "best plan", fontsize=8, color=ORANGE, rotation=90, va="center")
arrow(11, 14, 11, 47, color=TEAL)    # calibrate/realized -> forecaster
ax.text(8.5, 30.5, "realized data", fontsize=8, color=TEAL, rotation=90, va="center")
box(54, 1.5, 19, 7.5, "Real Data Center\n(System Data)", "#444444", fs=9)
arrow(63, 14, 63, 9, color=GREEN); arrow(70, 9, 90, 14, color=GREY)
ax.text(63, 11.4, "commands", ha="center", fontsize=7.5, color=GREEN)
fig.tight_layout(); fig.savefig(ASSETS / "fig3_arch.png", dpi=150); plt.close(fig)

# ---------------------------------------------------------------- Fig 4: validation result
flow = [4.8, 7.0, 9.3, 11.5, 13.8]
energy = [262579, 270195, 279772, 290474, 303267]
inlet = [23.00, 23.02, 23.85, 24.51, 25.22]
fig, ax1 = plt.subplots(figsize=(11, 5.6))
bars = ax1.bar(range(len(flow)), [e / 1000 for e in energy], color=LIGHT, edgecolor=NAVY,
               width=0.55, zorder=2)
bars[0].set_facecolor(ORANGE)
ax1.set_xticks(range(len(flow))); ax1.set_xticklabels([f"{f}" for f in flow])
ax1.set_xlabel("CRAH airflow setpoint (kg/s)", fontsize=10)
ax1.set_ylabel("hall cooling energy (MWh / 2 days)", fontsize=10, color=NAVY)
ax1.set_ylim(250, 310)
ax1.text(0, energy[0] / 1000 + 1.5, "optimum\n4.8 kg/s", ha="center", fontsize=9,
         color=ORANGE, weight="bold")
ax2 = ax1.twinx()
ax2.plot(range(len(flow)), inlet, "o-", color=RED, lw=2.2, zorder=3, label="peak server inlet")
ax2.axhline(26, color=RED, ls="--", lw=1.3); ax2.text(4.05, 26.1, "26 °C safety cap", color=RED, fontsize=9)
ax2.set_ylabel("peak server inlet temp (°C)", fontsize=10, color=RED)
ax2.set_ylim(22, 27)
ax1.set_title("Well-posed objective: cooling energy now varies 15.5% with airflow\n"
              "(was 0.07% before the fix) — search finds the cheapest safe point",
              fontsize=12, weight="bold", color=NAVY)
ax1.spines[["top"]].set_visible(False)
fig.tight_layout(); fig.savefig(ASSETS / "fig4_results.png", dpi=150); plt.close(fig)

print("figures rendered:", [p.name for p in sorted(ASSETS.glob("*.png"))])

# ============================================================ BUILD PPTX
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY_C = RGBColor(0x1F, 0x3A, 0x5F); ORANGE_C = RGBColor(0xE8, 0x85, 0x2A)
TEAL_C = RGBColor(0x2A, 0x8E, 0x9E); GREY_C = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); GREEN_C = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_C = RGBColor(0xEF, 0xF3, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def _txt(slide, x, y, w, h, text, size, color=NAVY_C, bold=False, align=PP_ALIGN.LEFT,
         italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = font
    return tb


def _rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def header(slide, title, kicker=None):
    _rect(slide, 0, 0, SW, Inches(1.15), NAVY_C)
    _rect(slide, 0, Inches(1.15), SW, Inches(0.06), ORANGE_C)
    _txt(slide, Inches(0.5), Inches(0.2), SW - Inches(1), Inches(0.8), title, 26, WHITE, True)
    if kicker:
        _txt(slide, Inches(0.55), Inches(0.82), SW - Inches(1), Inches(0.3), kicker, 12,
             RGBColor(0xBF, 0xD3, 0xE6), False)


def bullets(slide, items, x=Inches(0.7), y=Inches(1.5), w=None, size=16, gap=1.18):
    w = w or (SW - Inches(1.4))
    tb = slide.shapes.add_textbox(x, y, w, SH - y - Inches(0.4)); tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        lvl, text, col, bold = it if isinstance(it, tuple) else (0, it, GREY_C, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap; p.level = lvl
        r = p.add_run(); bullet = ("•  " if lvl == 0 else "–  ")
        r.text = bullet + text; f = r.font
        f.size = Pt(size - lvl * 2); f.color.rgb = col; f.bold = bold; f.name = "Calibri"
    return tb


def pic_slide(title, img, caption, kicker=None):
    s = prs.slides.add_slide(blank); header(s, title, kicker)
    from PIL import Image
    iw, ih = Image.open(img).size
    maxw, maxh = SW - Inches(1.2), SH - Inches(2.2)
    scale = min(maxw / iw, maxh / ih)
    w = Emu(int(iw * scale)); h = Emu(int(ih * scale))
    x = Emu(int((SW - w) / 2)); y = Inches(1.45)
    s.shapes.add_picture(str(img), x, y, w, h)
    _txt(s, Inches(0.6), SH - Inches(0.55), SW - Inches(1.2), Inches(0.45), caption, 12,
         GREY_C, False, PP_ALIGN.CENTER, italic=True)
    return s

# ---- Slide 1: title
s = prs.slides.add_slide(blank)
_rect(s, 0, 0, SW, SH, NAVY_C)
_rect(s, 0, Inches(4.05), SW, Inches(0.07), ORANGE_C)
_txt(s, Inches(0.9), Inches(2.1), SW - Inches(1.8), Inches(1.3),
     "Digital Twin Dual-Loop Control Framework", 40, WHITE, True)
_txt(s, Inches(0.95), Inches(3.25), SW - Inches(1.8), Inches(0.7),
     "Physics-based weekly cooling optimization for data centers", 20,
     RGBColor(0xE8, 0x85, 0x2A), False)
_txt(s, Inches(0.95), Inches(4.3), SW - Inches(1.8), Inches(0.6),
     "EnergyPlus 9.5 digital twin  ·  heuristic planner  ·  expert-supervised deployment", 15,
     RGBColor(0xBF, 0xD3, 0xE6))
_txt(s, Inches(0.95), Inches(5.6), SW - Inches(1.8), Inches(0.5),
     "3rd-draft progress report", 16, WHITE, True)
_txt(s, Inches(0.95), Inches(6.1), SW - Inches(1.8), Inches(0.4), "June 2026", 13,
     RGBColor(0xBF, 0xD3, 0xE6))

# ---- Slide 2: problem & goal
s = prs.slides.add_slide(blank); header(s, "The problem & the goal", "Why this project")
bullets(s, [
    (0, "Cooling is one of the largest energy costs in a data center; operators run conservative, fixed setpoints to stay safe — which wastes energy.", GREY_C, False),
    (0, "Goal: recommend, each week, the cooling setpoints that MINIMIZE energy…", NAVY_C, True),
    (1, "…while NEVER letting any server's inlet air exceed the 26 °C safety limit.", RED if False else GREY_C, False),
    (0, "Approach: a physics-based digital twin + a heuristic planner + an expert-supervised web app for deployment.", GREY_C, False),
    (0, "Three design commitments:", NAVY_C, True),
    (1, "Physics-based (first-principles EnergyPlus model — not a black box / surrogate).", GREY_C, False),
    (1, "Safe-by-construction (a hard temperature constraint, enforced three ways).", GREY_C, False),
    (1, "Architected to connect to a real data center (a single, well-defined deployment seam).", GREY_C, False),
], size=16)

# ---- Slide 3: architecture
pic_slide("How it works: the dual loop", ASSETS / "fig3_arch.png",
          "Inner loop searches safe setpoints in simulation; outer loop validates, deploys (via the BMS seam), and learns from realized data.",
          "Architecture")

# ---- Slide 4: physics twin
s = prs.slides.add_slide(blank); header(s, "A physics-based digital twin", "Physical fidelity")
bullets(s, [
    (0, "The twin is EnergyPlus 9.5 — the industry-standard building-physics engine — run live in Docker and coupled step-by-step to the planner (BCVTB).", GREY_C, False),
    (0, "First-principles, not a black box:", NAVY_C, True),
    (1, "Real building geometry + HVAC: chilled-water plant, air handlers, airflow, heat transfer.", GREY_C, False),
    (1, "Predicts per-rack inlet temperatures, zone temperatures, humidity, and power.", GREY_C, False),
    (0, "Three global setpoints the planner controls:", NAVY_C, True),
    (1, "CRAH supply-air temperature (20–26 °C) · airflow (4.8–13.8 kg/s) · chilled-water temp (13–19 °C).", GREY_C, False),
    (1, "Automatically broadcast to 45 individual actuators in the controlled hall.", GREY_C, False),
    (0, "Every candidate plan is scored by a FULL-WEEK EnergyPlus simulation (parallelized across CPU processes).", TEAL_C, True),
], size=15.5)

# ---- Slide 5: building topology
pic_slide("Data center topology — the building", ASSETS / "fig1_building.png",
          "Reconstructed from the real EnergyPlus zone geometry: 7 stacked halls; the planner controls hall 1F 2A; all halls share one chilled-water plant.",
          "Physical model")

# ---- Slide 6: hall detail
pic_slide("Data center topology — the controlled hall", ASSETS / "fig2_hall.png",
          "Hall 1F 2A: 22 agent-controlled CRAHs along the walls push cold air to the server inlets; the chilled-water plant feeds the coils. Inlet temp is the safety-critical signal.",
          "Physical model")

# ---- Slide 7: planner + safety
s = prs.slides.add_slide(blank); header(s, "The planner & safety nets", "Optimization under a hard constraint")
bullets(s, [
    (0, "Coarse-to-fine beam search over the 3 setpoints — lay a grid, keep the best few, zoom in, repeat.", GREY_C, False),
    (0, "Objective = minimize cooling energy + small comfort penalties (temperature, humidity).", GREY_C, False),
    (0, "Hard safety constraint: any week with a single step above 26 °C inlet is rejected outright.", RED if False else NAVY_C, True),
    (0, "Three independent safety nets:", NAVY_C, True),
    (1, "Hard cap during search (feasibility gate).", GREY_C, False),
    (1, "Robust re-check: re-simulate finalists on degraded 'weaker-than-expected plant' scenarios; keep only those safe in the worst case.", GREY_C, False),
    (1, "Deploy backstop: a breach on the deployed week is flagged 'deploy-blocked', never 'deployed'.", GREY_C, False),
    (0, "Optional day/night scheduling: different setpoints for day vs. night.", GREY_C, False),
], size=15.5)

# ---- Slide 8: web app
s = prs.slides.add_slide(blank); header(s, "The operator's cockpit (web app)", "Expert supervision")
bullets(s, [
    (0, "Full-stack web app (FastAPI + React) — single sign-on with operator vs. expert roles.", GREY_C, False),
    (0, "Six pages:", NAVY_C, True),
    (1, "New Plan — set the week + search settings, watch the search progress LIVE (streaming).", GREY_C, False),
    (1, "Review — recommended setpoints, energy vs. an as-operated baseline, inlet-temperature safety trajectory.", GREY_C, False),
    (1, "Digital Twin 3-D — interactive 3-D view of the hall with airflow.", GREY_C, False),
    (1, "History / Dashboard — trends of predicted vs. realized energy across weeks.", GREY_C, False),
    (0, "Expert workflow: approve · reject · edit setpoints · deploy — with the safety gate enforced (unsafe plans cannot be force-approved).", TEAL_C, True),
], size=15.5)

# ---- Slide 9: results
pic_slide("Result: the recommendations are physically well-posed", ASSETS / "fig4_results.png",
          "Verified with real EnergyPlus: cooling energy responds strongly + monotonically to airflow; the search picks the cheapest point that stays under the 26 °C cap.",
          "Validation")

# ---- Slide 10: status scorecard
s = prs.slides.add_slide(blank); header(s, "Status after the 3rd draft", "What is built")
bullets(s, [
    (0, "BUILT & TESTED (✓):  forecaster · EnergyPlus twin · beam-search planner · 3 safety nets · pre-validation · expert approve/edit/deploy workflow · 6-page web app.", GREEN_C, True),
    (1, "287 Python unit tests + 78 frontend tests pass; 8 Docker integration tests exercise the real EnergyPlus loop end-to-end.", GREY_C, False),
    (0, "PARTIAL (◐):  weather uses a fixed historical file (not forecast) · live telemetry monitoring · active humidity control · physics re-calibration.", ORANGE_C, True),
    (0, "NOT YET (○):  the connection to a real data center — today 'deploy' runs against a perturbed simulation, by design for this stage.", GREY_C, True),
    (0, "Honest fidelity caveats: deployment is simulation-to-simulation; the recirculation parameter is currently inert (safety margin is conservative/optimistic); numbers are directionally sound, not yet field-calibrated.", GREY_C, False),
], size=15)

# ---- Slide 11: ready to connect
s = prs.slides.add_slide(blank); header(s, "Ready to connect to a real data center", "The path to closed-loop")
bullets(s, [
    (0, "The deployment path is a single, well-defined seam:  BmsAdapter.apply(setpoints, week).", NAVY_C, True),
    (0, "Everything downstream of that seam is already built and tested:", NAVY_C, True),
    (1, "expert approval gate · realized-KPI capture · automatic calibration ingest · deploy-blocked safety backstop.", GREY_C, False),
    (0, "Connect-up work (well-scoped):", NAVY_C, True),
    (1, "1) Implement the BMS adapter (BACnet / Modbus / vendor API) — start in shadow mode (recommend-only, no actuation).", GREY_C, False),
    (1, "2) Ingest live telemetry (inlet temps, power, humidity) into the calibration loop + a live monitoring dashboard.", GREY_C, False),
    (1, "3) Validate predicted vs. truly realized for several weeks, THEN enable closed-loop actuation.", GREY_C, False),
    (0, "Design choices that make this clean: physics-based twin (transfers to the real site), strict separation of 'recommend' vs. 'actuate', and a calibration loop already built to absorb real data.", TEAL_C, True),
], size=14.5)

# ---- Slide 12: roadmap
s = prs.slides.add_slide(blank); header(s, "Roadmap", "Next steps, in priority order")
bullets(s, [
    (0, "Tier A — Close the loop to reality (highest impact):", NAVY_C, True),
    (1, "BMS adapter (shadow mode first) · real telemetry ingest · live monitoring + alerts dashboard.", GREY_C, False),
    (0, "Tier B — Raise twin fidelity:", NAVY_C, True),
    (1, "Fix the recirculation physics, then calibrate to measured rack inlets · activate physics re-calibration · per-rack drill-down.", GREY_C, False),
    (0, "Tier C — Smarter planning:", NAVY_C, True),
    (1, "Weather forecasting + uncertainty · carbon / time-of-use-tariff-aware objective.", GREY_C, False),
    (0, "Tier D — Operator experience & assurance:", NAVY_C, True),
    (1, "Mid-week override UI · failure explainability · end-to-end tests · a recommend-only shadow pilot.", GREY_C, False),
], size=15.5)

# ---- Slide 13: summary
s = prs.slides.add_slide(blank); header(s, "Summary", "Where the 3rd draft stands")
bullets(s, [
    (0, "A physics-based digital twin (EnergyPlus) + a heuristic planner + an expert-supervised web app — validated end-to-end in simulation.", GREY_C, False),
    (0, "Safe-by-construction: a hard 26 °C inlet limit enforced by three independent nets.", GREY_C, False),
    (0, "Recommendations are physically well-posed and directionally sound (e.g. +3.2% energy saving with zero violations); not yet field-calibrated.", GREY_C, False),
    (0, "Architected and ready to connect to a real data center through one adapter seam — shadow mode first, then closed loop.", TEAL_C, True),
    (0, "Next: connect telemetry + BMS, raise fidelity (recirculation calibration), and run a shadow pilot.", NAVY_C, True),
], size=16)

pptx_path = OUT / "DCTwin_3rd_draft_report.pptx"
prs.save(str(pptx_path))
print("SAVED:", pptx_path, "slides:", len(prs.slides._sldIdLst))
